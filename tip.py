#!/usr/bin/env python
# -*- coding: UTF-8 -*-

from pptx import Presentation

_tpl = './tpl/law.pptx'
_output = './output/result.pptx'
_content = './content.txt'
SLD_LAYOUT_TITLE_AND_CONTENT = 0

prs = Presentation(_tpl)


def usage():
    pass


def read_content():
    """
    从content.txt中读取每一页的内容
    """
    slides = []
    current = {}
    with open(_content, 'r', encoding='utf8') as f:
        for l in f.readlines():
            l = l.strip()  # 把末尾的'\n'删掉
            if not l:
                continue
            if l.startswith('#'):  # 开始下一页
                if current:
                    if(len(current['paragraphs']) > 0):
                        slides.append(current)
                    current = {}

            elif not 'title' in current:
                current['title'] = l

            else:
                if not 'paragraphs' in current:
                    current['paragraphs'] = [l]
                else:
                    current['paragraphs'].append(l)
        if current:
            slides.append(current)

    return slides


def gen():
    slide_contents = read_content()
    if len(slide_contents) == 0:
        print('content.txt中没有有效内容')
        return
    slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
    for content in slide_contents:
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes[0]
        title_shape.text_frame.paragraphs[0].text = content['title']

        body_shape = slide.shapes[1]
        body_shape.text_frame.paragraphs[0].text = content['paragraphs'][0]
        for p in content['paragraphs'][1:]:
            body_shape.text_frame.add_paragraph().text = p

    prs.save(_output)


if __name__ == '__main__':
    gen()
